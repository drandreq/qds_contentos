import io
import os
import logging
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from core.models import LessonMetadata, SlideDirective

logger = logging.getLogger(__name__)

class AssetFactory:
    """
    Sprint 8: Converts parsed MP-Dialect JSON models into tangible visual assets.
    Generates single PNG slides or entire PPTX decks.
    """

    def generate_png(self, slide: SlideDirective) -> bytes:
        """
        Creates a simple 1920x1080 PNG image of the slide content.
        Uses a dark gray background with centered white text.
        Returns the image as bytes.
        """
        width, height = 1920, 1080
        background_color = (40, 44, 52) # Dark gray (Dracula/OneDark vibe)
        text_color = (255, 255, 255)

        # Create image
        image = Image.new("RGB", (width, height), color=background_color)
        draw = ImageDraw.Draw(image)

        # Attempt to load a default TrueType font, fallback to standard if not available
        try:
            # Arial is usually available on Windows and some Linux. 
            # In a tiny Alpine/Slim Docker, we might need a specific path or just rely on default.
            font = ImageFont.truetype("arial.ttf", 80)
        except Exception:
            try:
                # Fallback for some Linux 
                font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 80)
            except Exception:
                # Absolute fallback (bitmap, cannot set size easily)
                font = ImageFont.load_default()

        # Handle text wrapping manually or just center as best as possible for now
        # For a robust solution, textwrap module is ideal, but let's do a simple bounded box.
        text = slide.content
        
        # Calculate text bounding box
        bbox = draw.multiline_textbbox((0, 0), text, font=font, align="center")
        text_width = bbox[2] - bbox[0]
        text_height = bbox[3] - bbox[1]

        # Draw centered
        x = (width - text_width) / 2
        y = (height - text_height) / 2
        
        draw.multiline_text((x, y), text, fill=text_color, font=font, align="center")

        # Save to memory
        img_byte_arr = io.BytesIO()
        image.save(img_byte_arr, format='PNG')
        return img_byte_arr.getvalue()

    def generate_pptx(self, metadata: LessonMetadata) -> bytes:
        """
        Generates a complete PPTX presentation from a LessonMetadata model.
        Returns the .pptx file as bytes.
        """
        prs = Presentation()

        for slide_data in metadata.slides:
            # For simplicity, if layout says 'Title', use the Title layout (0).
            # Otherwise use the Title and Content layout (1).
            layout_idx = 0 if "title" in slide_data.layout.lower() else 1
            
            # Avoid out of bounds if layout doesn't exist
            if layout_idx >= len(prs.slide_layouts):
                layout_idx = 0
                
            slide_layout = prs.slide_layouts[layout_idx]
            pres_slide = prs.slides.add_slide(slide_layout)
            
            # Populate content depending on the layout type
            if layout_idx == 0:
                # Title layout
                title_shape = pres_slide.shapes.title
                title_shape.text = slide_data.content
            else:
                # Title & Content Layout
                # We put the text in the main body (idx 1 usually)
                try:
                    title_shape = pres_slide.shapes.title
                    title_shape.text = "ContentOS Auto-Slide" # Generic title
                    
                    body_shape = pres_slide.shapes.placeholders[1]
                    tf = body_shape.text_frame
                    tf.text = slide_data.content
                except Exception as e:
                    logger.warning(f"Could not format slide cleanly: {e}")

        # Save to memory bytes
        pptx_bytes = io.BytesIO()
        prs.save(pptx_bytes)
        return pptx_bytes.getvalue()

# Singleton
asset_factory = AssetFactory()
